#!/usr/bin/env python3
"""Generate A Piece of Peace Business & Strategic Plan PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Soft light palette
DARK_BG = RGBColor(0xF8, 0xF6, 0xF2)       # warm off-white
PURPLE = RGBColor(0x9B, 0x8B, 0xC4)         # soft lavender
TEAL = RGBColor(0x7E, 0xB5, 0xA6)           # muted sage
GOLD = RGBColor(0xD4, 0xB4, 0x83)           # warm sand
WHITE = RGBColor(0x3A, 0x36, 0x3B)          # dark charcoal (headings)
MUTED = RGBColor(0x6B, 0x66, 0x6A)          # warm gray (body text)
SURFACE = RGBColor(0xEE, 0xEB, 0xE5)        # light warm card bg
ROSE = RGBColor(0xD4, 0x8B, 0x9A)           # soft blush
SKY = RGBColor(0x89, 0xAF, 0xC8)            # dusty blue

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=DARK_BG):
    """Fill slide background with color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Helvetica Neue"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=MUTED, bullet_color=TEAL):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25cf  "
        run_bullet.font.size = Pt(font_size - 2)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Helvetica Neue"
        # Item text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Helvetica Neue"
    return tf

def add_accent_line(slide, left, top, width, color=PURPLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_card(slide, left, top, width, height, color=SURFACE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(0xDE, 0xDA, 0xD3)
    shape.line.width = Pt(1)
    return shape

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)
add_accent_line(slide, 1, 2.2, 2.5, PURPLE)
add_textbox(slide, 1, 2.5, 11, 1.2, "A Piece of Peace", 54, PURPLE, True)
add_textbox(slide, 1, 3.5, 11, 0.8,
            "A Sound-Based Wellness App", 28, TEAL)
add_textbox(slide, 1, 4.4, 11, 0.6,
            "Business & Strategic Plan", 20, MUTED)
add_textbox(slide, 1, 5.6, 11, 0.5,
            "Presented by Rhonda Davis  \u2022  Apple Developer Academy", 16, MUTED)

# ============================================================
# SLIDE 2: Executive Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, TEAL)
add_textbox(slide, 1, 1.0, 11, 0.7, "Executive Summary", 36, WHITE, True)

add_textbox(slide, 1, 2.0, 10.5, 1.0,
            "A Piece of Peace is a mobile wellness app that uses interactive, real-time sound "
            "as its primary therapeutic interface. Unlike conventional wellness apps that "
            "rely on journaling, tracking, or social features, A Piece of Peace puts users in "
            "direct relationship with music and sound to regulate their emotional and "
            "physiological state.",
            17, MUTED)

add_card(slide, 1, 3.5, 3.3, 2.8)
add_textbox(slide, 1.3, 3.7, 2.8, 0.4, "Mission", 14, TEAL, True)
add_textbox(slide, 1.3, 4.1, 2.8, 1.8,
            "Make emotional regulation accessible through the universal language of sound.",
            16, MUTED)

add_card(slide, 4.6, 3.5, 3.3, 2.8)
add_textbox(slide, 4.9, 3.7, 2.8, 0.4, "Vision", 14, PURPLE, True)
add_textbox(slide, 4.9, 4.1, 2.8, 1.8,
            "A world where everyone has an intuitive, non-verbal tool for wellness \u2014 "
            "as natural as pressing play.",
            16, MUTED)

add_card(slide, 8.2, 3.5, 3.3, 2.8)
add_textbox(slide, 8.5, 3.7, 2.8, 0.4, "Core Differentiator", 14, GOLD, True)
add_textbox(slide, 8.5, 4.1, 2.8, 1.8,
            "Sound is generated in real-time, not pre-recorded. "
            "The app responds to the user's body, breath, and rhythm.",
            16, MUTED)

# ============================================================
# SLIDE 3: Problem Statement
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, ROSE)
add_textbox(slide, 1, 1.0, 11, 0.7, "The Problem", 36, WHITE, True)

problems = [
    "76% of U.S. adults report stress-related health impacts (APA, 2023)",
    "Most wellness apps rely on text-heavy interfaces \u2014 journaling, CBT worksheets, tracking dashboards",
    "These approaches create cognitive load for people who are already overwhelmed",
    "Music is proven to regulate heart rate, cortisol, and mood \u2014 but no app makes it interactive",
    "Gen Z and millennials want sensory, immersive experiences \u2014 not more screen reading",
]
add_bullet_list(slide, 1, 2.0, 10.5, 4.5, problems, 17, MUTED, ROSE)

# ============================================================
# SLIDE 4: Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, TEAL)
add_textbox(slide, 1, 1.0, 11, 0.7, "Our Solution", 36, WHITE, True)

add_textbox(slide, 1, 2.0, 10.5, 0.8,
            "A Piece of Peace replaces passive content consumption with active sound interaction.",
            18, WHITE)

features = [
    ("Guided Breathing + Live Audio", "Breathe in sync with ambient soundscapes generated in real-time at 432Hz"),
    ("Mood Mapping", "Translate emotions into sound using sliders, colors, and words \u2014 the app finds your match"),
    ("Tap to Regulate", "Tap your natural rhythm \u2014 the app detects your BPM and shapes calming sound around it"),
    ("Hydration & Movement", "Gentle stretch guidance and water reminders paired with adaptive music"),
]

for i, (title, desc) in enumerate(features):
    y = 3.0 + i * 1.1
    add_card(slide, 1, y, 10.5, 0.9)
    add_textbox(slide, 1.3, y + 0.1, 3, 0.4, title, 16, TEAL, True)
    add_textbox(slide, 4.5, y + 0.15, 6.8, 0.6, desc, 15, MUTED)

# ============================================================
# SLIDE 5: Target Market
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, PURPLE)
add_textbox(slide, 1, 1.0, 11, 0.7, "Target Market", 36, WHITE, True)

add_card(slide, 1, 2.2, 3.5, 3.5)
add_textbox(slide, 1.3, 2.4, 3, 0.4, "Primary", 14, PURPLE, True)
add_textbox(slide, 1.3, 2.8, 3, 0.4, "Ages 18\u201334", 22, WHITE, True)
add_bullet_list(slide, 1.3, 3.3, 3, 2, [
    "College students managing stress",
    "Young professionals with anxiety",
    "Music-oriented, sensory learners",
    "Digital natives seeking non-verbal tools",
], 14, MUTED, PURPLE)

add_card(slide, 5, 2.2, 3.5, 3.5)
add_textbox(slide, 5.3, 2.4, 3, 0.4, "Secondary", 14, TEAL, True)
add_textbox(slide, 5.3, 2.8, 3, 0.4, "Ages 35\u201355", 22, WHITE, True)
add_bullet_list(slide, 5.3, 3.3, 3, 2, [
    "Wellness-conscious parents",
    "Therapists recommending tools",
    "Meditation practitioners",
    "People with sensory processing needs",
], 14, MUTED, TEAL)

add_card(slide, 9, 2.2, 3.5, 3.5)
add_textbox(slide, 9.3, 2.4, 3, 0.4, "Niche", 14, GOLD, True)
add_textbox(slide, 9.3, 2.8, 3, 0.4, "Specialized", 22, WHITE, True)
add_bullet_list(slide, 9.3, 3.3, 3, 2, [
    "Music therapy patients",
    "ADHD / neurodivergent users",
    "Trauma recovery support",
    "Sleep and insomnia management",
], 14, MUTED, GOLD)

# ============================================================
# SLIDE 6: Market Opportunity
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, GOLD)
add_textbox(slide, 1, 1.0, 11, 0.7, "Market Opportunity", 36, WHITE, True)

add_card(slide, 1, 2.2, 3.5, 2)
add_textbox(slide, 1.3, 2.4, 3, 0.4, "Global Wellness App Market", 13, GOLD, True)
add_textbox(slide, 1.3, 2.9, 3, 0.5, "$7.8B", 40, WHITE, True)
add_textbox(slide, 1.3, 3.5, 3, 0.4, "by 2027 (CAGR 17.7%)", 14, MUTED)

add_card(slide, 5, 2.2, 3.5, 2)
add_textbox(slide, 5.3, 2.4, 3, 0.4, "Mental Health App Market", 13, TEAL, True)
add_textbox(slide, 5.3, 2.9, 3, 0.5, "$5.2B", 40, WHITE, True)
add_textbox(slide, 5.3, 3.5, 3, 0.4, "by 2027 (CAGR 15.9%)", 14, MUTED)

add_card(slide, 9, 2.2, 3.5, 2)
add_textbox(slide, 9.3, 2.4, 3, 0.4, "Music Therapy Market", 13, PURPLE, True)
add_textbox(slide, 9.3, 2.9, 3, 0.5, "$2.1B", 40, WHITE, True)
add_textbox(slide, 9.3, 3.5, 3, 0.4, "by 2028 (CAGR 12.4%)", 14, MUTED)

trends = [
    "Post-pandemic mental health awareness driving sustained demand for accessible tools",
    "Gen Z preference for sensory/interactive experiences over text-based interfaces",
    "Apple Health and wellness ecosystem integration creating distribution opportunities",
    "Growing clinical evidence supporting music therapy and sound-based interventions",
]
add_bullet_list(slide, 1, 4.7, 10.5, 2.5, trends, 15, MUTED, GOLD)

# ============================================================
# SLIDE 7: Product Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, TEAL)
add_textbox(slide, 1, 1.0, 11, 0.7, "Product Overview", 36, WHITE, True)

screens = [
    ("Home", "\u2302", "Time-aware greeting, wellness overview, now playing, daily quote from API", PURPLE),
    ("Wellness", "\U0001F33F", "Guided breathing with live audio synthesis, stretch routines, hydration", TEAL),
    ("Music", "\u266B", "Mood mapping with sliders + color + words, tap-to-regulate with BPM detection", GOLD),
    ("Connect", "\U0001F91D", "Community features, friend check-ins, shared experiences", ROSE),
]

for i, (name, icon, desc, color) in enumerate(screens):
    x = 1 + i * 2.9
    add_card(slide, x, 2.2, 2.6, 4)
    add_textbox(slide, x + 0.2, 2.5, 2.2, 0.5, icon, 32, color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 3.1, 2.2, 0.4, name, 20, WHITE, True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 3.6, 2.2, 2.2, desc, 14, MUTED,
                alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8: Technology Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, SKY)
add_textbox(slide, 1, 1.0, 11, 0.7, "Technology Architecture", 36, WHITE, True)

tech_items = [
    ("SwiftUI", "Declarative UI framework \u2014 reactive state management with @Observable, "
     "@State, and @Bindable for real-time UI updates"),
    ("AVAudioEngine", "Core Audio framework for real-time sound synthesis \u2014 generates ambient tones, "
     "binaural beats, and bell chimes from raw sine waves on the audio thread"),
    ("Structured Concurrency", "Swift async/await for breathing timers, API calls, and task cancellation \u2014 "
     "no memory leaks or orphaned processes"),
    ("Custom Layout Protocol", "FlowLayout for responsive word-tag wrapping, similar to CSS flexbox"),
    ("REST API Integration", "Live daily quotes from ZenQuotes API with graceful fallback on failure"),
    ("Audio-Thread Safety", "Lock-free render state pattern using @unchecked Sendable for real-time "
     "audio without blocking the main thread"),
]

for i, (title, desc) in enumerate(tech_items):
    y = 2.0 + i * 0.85
    add_textbox(slide, 1, y, 2.5, 0.4, title, 16, SKY, True)
    add_textbox(slide, 3.8, y, 8, 0.7, desc, 14, MUTED)

# ============================================================
# SLIDE 9: Competitive Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, PURPLE)
add_textbox(slide, 1, 1.0, 11, 0.7, "Competitive Landscape", 36, WHITE, True)

# Table header
header_y = 2.2
cols = [("Feature", 2.5), ("Calm", 1.8), ("Headspace", 1.8), ("Endel", 1.8), ("A Piece of Peace", 2)]
x = 1
for label, w in cols:
    c = GOLD if label == "A Piece of Peace" else MUTED
    add_textbox(slide, x, header_y, w, 0.4, label, 14, c, True, PP_ALIGN.CENTER)
    x += w

rows = [
    ("Real-time audio synthesis", "\u2717", "\u2717", "\u2713", "\u2713"),
    ("Breath-synced sound", "\u2717", "\u25CB", "\u25CB", "\u2713"),
    ("Tap rhythm detection", "\u2717", "\u2717", "\u2717", "\u2713"),
    ("Mood-to-music mapping", "\u2717", "\u2717", "\u25CB", "\u2713"),
    ("No subscription required*", "\u2717", "\u2717", "\u2717", "\u2713"),
    ("Binaural beat generation", "\u2717", "\u2717", "\u2713", "\u2713"),
    ("Interactive (not passive)", "\u2717", "\u25CB", "\u25CB", "\u2713"),
]

for i, (feature, *vals) in enumerate(rows):
    y = 2.7 + i * 0.55
    bg_color = RGBColor(0xF2, 0xEF, 0xE9) if i % 2 == 0 else SURFACE
    add_card(slide, 1, y, 10.1, 0.48, bg_color)
    add_textbox(slide, 1.2, y + 0.05, 2.3, 0.38, feature, 13, MUTED)
    x = 3.5
    for j, v in enumerate(vals):
        c = GOLD if j == 3 and v == "\u2713" else (RGBColor(0x7E, 0xB5, 0x8E) if v == "\u2713" else (RGBColor(0xD4, 0xB4, 0x83) if v == "\u25CB" else RGBColor(0x9E, 0x99, 0x97)))
        add_textbox(slide, x, y + 0.05, 1.8, 0.38, v, 16, c, alignment=PP_ALIGN.CENTER)
        x += 1.8

add_textbox(slide, 1, 6.6, 10, 0.4,
            "\u2713 = Full support   \u25CB = Partial   \u2717 = Not available   *Launch strategy",
            12, RGBColor(0x9E, 0x99, 0x97))

# ============================================================
# SLIDE 10: Business Model
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, GOLD)
add_textbox(slide, 1, 1.0, 11, 0.7, "Business Model", 36, WHITE, True)

add_card(slide, 1, 2.2, 3.5, 3.8)
add_textbox(slide, 1.3, 2.4, 3, 0.4, "Free Tier", 16, RGBColor(0x7E, 0xB5, 0x8E), True)
add_textbox(slide, 1.3, 2.8, 3, 0.4, "$0", 32, WHITE, True)
add_bullet_list(slide, 1.3, 3.3, 3, 2.5, [
    "Guided breathing (1 mode)",
    "Mood mapping",
    "Tap to regulate",
    "Daily quote",
    "Basic stretch routine",
], 13, MUTED, RGBColor(0x7E, 0xB5, 0x8E))

add_card(slide, 5, 2.2, 3.5, 3.8)
add_textbox(slide, 5.3, 2.4, 3, 0.4, "Premium", 16, PURPLE, True)
add_textbox(slide, 5.3, 2.8, 3, 0.4, "$4.99/mo", 32, WHITE, True)
add_bullet_list(slide, 5.3, 3.3, 3, 2.5, [
    "All 3 breathing soundscapes",
    "Binaural beat modes",
    "Extended stretch library",
    "Personalized mood history",
    "Offline sound generation",
], 13, MUTED, PURPLE)

add_card(slide, 9, 2.2, 3.5, 3.8)
add_textbox(slide, 9.3, 2.4, 3, 0.4, "Annual", 16, GOLD, True)
add_textbox(slide, 9.3, 2.8, 3, 0.4, "$39.99/yr", 32, WHITE, True)
add_bullet_list(slide, 9.3, 3.3, 3, 2.5, [
    "Everything in Premium",
    "Priority new features",
    "Family sharing (up to 5)",
    "Custom soundscape builder",
    "Save 33% vs monthly",
], 13, MUTED, GOLD)

add_textbox(slide, 1, 6.3, 10.5, 0.5,
            "Additional revenue: B2B licensing for therapy practices, corporate wellness programs, "
            "and educational institutions",
            14, MUTED)

# ============================================================
# SLIDE 11: Go-to-Market Strategy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, TEAL)
add_textbox(slide, 1, 1.0, 11, 0.7, "Go-to-Market Strategy", 36, WHITE, True)

phases = [
    ("Phase 1: Soft Launch", "Months 1\u20133", TEAL, [
        "Launch on iOS App Store (free tier)",
        "Target university wellness centers",
        "Apple Developer Academy network",
        "Beta testing with 200\u2013500 users",
    ]),
    ("Phase 2: Growth", "Months 4\u20139", PURPLE, [
        "Introduce Premium subscription",
        "Influencer partnerships (wellness/music)",
        "App Store Optimization (ASO)",
        "Content marketing: sound + science blog",
    ]),
    ("Phase 3: Scale", "Months 10\u201318", GOLD, [
        "B2B partnerships (therapists, corporates)",
        "Android expansion",
        "Apple Watch companion app",
        "International localization",
    ]),
]

for i, (title, timeline, color, items) in enumerate(phases):
    x = 1 + i * 3.7
    add_card(slide, x, 2.2, 3.4, 4.5)
    add_textbox(slide, x + 0.2, 2.4, 3, 0.4, title, 16, color, True)
    add_textbox(slide, x + 0.2, 2.8, 3, 0.3, timeline, 13, MUTED)
    add_bullet_list(slide, x + 0.2, 3.3, 3, 3, items, 14, MUTED, color)

# ============================================================
# SLIDE 12: Marketing Strategy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, ROSE)
add_textbox(slide, 1, 1.0, 11, 0.7, "Marketing Strategy", 36, WHITE, True)

channels = [
    ("\U0001F3B5 Content Marketing",
     "\"The Science of Sound\" blog series, YouTube demos of real-time sound generation, "
     "podcast guest appearances on wellness and music production shows"),
    ("\U0001F4F1 Social Media",
     "Short-form video (TikTok/Reels) showing mood-to-music transformation, "
     "ASMR-style breathing demos, behind-the-scenes of audio synthesis"),
    ("\U0001F3EB Institutional Partnerships",
     "University wellness centers, therapy practices, corporate wellness programs. "
     "Offer free pilot programs with usage analytics"),
    ("\U0001F34E Apple Ecosystem",
     "Target Apple Today tab features, leverage Developer Academy network, "
     "pursue Apple Design Award recognition for innovative use of AVAudioEngine"),
]

for i, (title, desc) in enumerate(channels):
    y = 2.0 + i * 1.3
    add_card(slide, 1, y, 10.5, 1.1)
    add_textbox(slide, 1.3, y + 0.1, 3, 0.4, title, 16, ROSE, True)
    add_textbox(slide, 1.3, y + 0.5, 9.8, 0.5, desc, 14, MUTED)

# ============================================================
# SLIDE 13: Financial Projections
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, GOLD)
add_textbox(slide, 1, 1.0, 11, 0.7, "Financial Projections", 36, WHITE, True)

# Year cards
years = [
    ("Year 1", "5,000", "500", "$25K", "Build & validate"),
    ("Year 2", "25,000", "3,750", "$180K", "Growth & retention"),
    ("Year 3", "100,000", "20,000", "$960K", "Scale & B2B"),
]

for i, (year, users, premium, revenue, focus) in enumerate(years):
    x = 1 + i * 3.7
    add_card(slide, x, 2.2, 3.4, 3.5)
    add_textbox(slide, x + 0.2, 2.4, 3, 0.4, year, 18, GOLD, True)
    add_textbox(slide, x + 0.2, 2.9, 3, 0.3, "Total Users", 12, MUTED)
    add_textbox(slide, x + 0.2, 3.2, 3, 0.4, users, 28, WHITE, True)
    add_textbox(slide, x + 0.2, 3.7, 3, 0.3, "Premium Subscribers", 12, MUTED)
    add_textbox(slide, x + 0.2, 4.0, 3, 0.4, premium, 22, TEAL, True)
    add_textbox(slide, x + 0.2, 4.5, 3, 0.3, f"Projected ARR: {revenue}", 14, GOLD)
    add_textbox(slide, x + 0.2, 4.9, 3, 0.3, f"Focus: {focus}", 12, MUTED)

assumptions = [
    "Assumes 10% free-to-premium conversion (industry avg: 5\u201312%)",
    "Premium ARPU of $48/year (blended monthly + annual)",
    "B2B revenue not included in Year 1\u20132 projections",
    "CAC target: <$3.00 through organic and partnership channels",
]
add_bullet_list(slide, 1, 6.0, 10.5, 1.5, assumptions, 13, RGBColor(0x9E, 0x99, 0x97), GOLD)

# ============================================================
# SLIDE 14: Development Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, PURPLE)
add_textbox(slide, 1, 1.0, 11, 0.7, "Development Roadmap", 36, WHITE, True)

milestones = [
    ("Q2 2026", "MVP Launch", [
        "Core 4-screen app on App Store",
        "3 breathing soundscapes",
        "Mood mapping + tap-to-regulate",
        "Daily quote API integration",
    ], TEAL),
    ("Q3 2026", "Premium Features", [
        "Subscription with StoreKit 2",
        "Expanded sound library",
        "Offline audio generation",
        "Usage analytics integration",
    ], PURPLE),
    ("Q4 2026", "Platform Expansion", [
        "Apple Watch haptic breathing",
        "HealthKit integration",
        "Siri Shortcuts support",
        "Widget for Home Screen",
    ], GOLD),
    ("Q1 2027", "Scale & Ecosystem", [
        "Android version (Kotlin)",
        "B2B dashboard for therapists",
        "API for third-party integration",
        "AI-driven mood prediction",
    ], ROSE),
]

for i, (quarter, title, items, color) in enumerate(milestones):
    x = 1 + i * 2.9
    add_card(slide, x, 2.2, 2.6, 4.5)
    add_textbox(slide, x + 0.2, 2.4, 2.2, 0.3, quarter, 13, color, True,
                PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 2.7, 2.2, 0.4, title, 17, WHITE, True,
                PP_ALIGN.CENTER)
    add_bullet_list(slide, x + 0.2, 3.3, 2.2, 3, items, 13, MUTED, color)

# ============================================================
# SLIDE 15: Risk Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, ROSE)
add_textbox(slide, 1, 1.0, 11, 0.7, "Risk Analysis & Mitigation", 36, WHITE, True)

risks = [
    ("Market Saturation",
     "Crowded wellness app space with established players",
     "Differentiate through interactive sound (no competitor offers real-time synthesis + rhythm detection)"),
    ("User Retention",
     "Wellness apps average 4% 30-day retention",
     "Sound is inherently repeatable; adaptive audio creates unique sessions each time"),
    ("Technical Complexity",
     "Real-time audio on mobile requires careful performance optimization",
     "AVAudioEngine is Apple's native framework; audio-thread architecture already proven in v1"),
    ("Monetization",
     "Users resistant to wellness app subscriptions",
     "Generous free tier builds habit; premium unlocks depth, not access"),
]

for i, (risk, desc, mitigation) in enumerate(risks):
    y = 2.0 + i * 1.3
    add_card(slide, 1, y, 10.5, 1.1)
    add_textbox(slide, 1.3, y + 0.05, 2.5, 0.35, risk, 15, ROSE, True)
    add_textbox(slide, 1.3, y + 0.4, 4.5, 0.35, desc, 13, MUTED)
    add_textbox(slide, 6.0, y + 0.05, 1.2, 0.35, "Mitigation:", 12, TEAL, True)
    add_textbox(slide, 6.0, y + 0.4, 5.2, 0.65, mitigation, 13, MUTED)

# ============================================================
# SLIDE 16: Key Metrics / KPIs
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, TEAL)
add_textbox(slide, 1, 1.0, 11, 0.7, "Key Performance Indicators", 36, WHITE, True)

kpis = [
    ("DAU / MAU Ratio", "> 25%", "Measures daily engagement stickiness", SKY),
    ("Session Duration", "> 4 min", "Breathing session = 4\u20138 min; indicates full use", TEAL),
    ("Free \u2192 Premium", "> 8%", "Conversion rate from free to paid tier", PURPLE),
    ("30-Day Retention", "> 15%", "Industry avg is 4%; sound-based habit is stickier", GOLD),
    ("NPS Score", "> 50", "Net Promoter Score from in-app surveys", RGBColor(0x7E, 0xB5, 0x8E)),
    ("Audio Sessions / User / Week", "> 3", "Core engagement metric \u2014 are users coming back to sound?", ROSE),
]

for i, (metric, target, desc, color) in enumerate(kpis):
    y = 2.1 + i * 0.85
    add_card(slide, 1, y, 10.5, 0.7)
    add_textbox(slide, 1.3, y + 0.15, 3, 0.4, metric, 15, WHITE, True)
    add_textbox(slide, 4.5, y + 0.15, 1.5, 0.4, target, 16, color, True, PP_ALIGN.CENTER)
    add_textbox(slide, 6.2, y + 0.15, 5, 0.4, desc, 14, MUTED)

# ============================================================
# SLIDE 17: Team & Resources Needed
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 0.8, 2, PURPLE)
add_textbox(slide, 1, 1.0, 11, 0.7, "Team & Resources", 36, WHITE, True)

roles = [
    ("iOS Developer", "SwiftUI + AVAudioEngine expertise", "In place"),
    ("UX/UI Designer", "Accessibility, dark-mode-first design system", "Needed"),
    ("Sound Designer", "Ambient soundscape creation, frequency tuning", "Needed"),
    ("Marketing Lead", "App Store, social media, partnerships", "Needed"),
    ("Backend Developer", "User accounts, analytics, subscription mgmt", "Phase 2"),
]

for i, (role, desc, status) in enumerate(roles):
    y = 2.2 + i * 0.9
    add_card(slide, 1, y, 10.5, 0.75)
    add_textbox(slide, 1.3, y + 0.15, 2.5, 0.4, role, 15, WHITE, True)
    add_textbox(slide, 4, y + 0.15, 5, 0.4, desc, 14, MUTED)
    sc = RGBColor(0x7E, 0xB5, 0x8E) if status == "In place" else (GOLD if status == "Phase 2" else ROSE)
    add_textbox(slide, 9.5, y + 0.15, 2, 0.4, status, 14, sc, True, PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.8, 10.5, 0.5,
            "Estimated seed funding needed: $50K\u2013$75K (covers 12 months: design, sound, marketing, App Store fees)",
            15, MUTED)

# ============================================================
# SLIDE 18: Call to Action
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_line(slide, 4.5, 2.0, 4.3, PURPLE)
add_textbox(slide, 1, 2.3, 11.3, 1.0, "Let\u2019s Build the Future\nof Sound-Based Wellness",
            40, WHITE, True, PP_ALIGN.CENTER)

add_textbox(slide, 2, 3.8, 9.3, 0.8,
            "A Piece of Peace isn\u2019t just an app \u2014 it\u2019s a new paradigm for how technology\n"
            "can support emotional regulation through the universal language of music.",
            18, MUTED, alignment=PP_ALIGN.CENTER)

add_card(slide, 3.5, 5.0, 6.3, 1.5)
add_textbox(slide, 3.8, 5.1, 5.7, 0.5, "Next Steps", 18, TEAL, True, PP_ALIGN.CENTER)
add_bullet_list(slide, 4.2, 5.5, 5, 1, [
    "Team review and feedback",
    "Finalize feature scope per Academy guidelines",
    "Begin user testing with cohort",
], 15, MUTED, TEAL)

add_textbox(slide, 1, 6.8, 11.3, 0.4, "A Piece of Peace  \u2022  Rhonda Davis  \u2022  Apple Developer Academy  \u2022  2026",
            14, RGBColor(0x9E, 0x99, 0x97), alignment=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = "/Users/rhondadavis/Library/Mobile Documents/com~apple~CloudDocs/Resonance/A_Piece_of_Peace_Business_Plan.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
